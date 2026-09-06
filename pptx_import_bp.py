"""Template-Import aus PPTX-Dateien und Canva-Designs (Phase B3).

Canva-Link oder PPTX rein, fertiges MemeOS-Template raus: Hintergrundbild (alle Bilder und
Farbflächen flach zusammengesetzt) plus pil_config nach Renderer-Vertrag (memeos_render).
Textboxen bleiben editierbare Text-Elemente; Platzhalter wie {problem_place} werden Variablen.

Blueprint 'pptx_import', Routen (Login):
  POST /api/templates/import/pptx    multipart 'file' + name/category/series/dry_run
  POST /api/templates/import/canva   JSON {design_url, name, category, series, dry_run}
  GET  /api/templates/import/help    Platzhalter-Anleitung + Variablenliste

Kein app-Import auf Modulebene (zirkulär) – innerhalb der Funktionen 'import app as appmod'.
"""
import io
import os
import re
import json
import time
import logging
from datetime import datetime
from functools import wraps

import requests
from flask import Blueprint, request, jsonify, session, redirect
from PIL import Image, ImageDraw

from models import db, MemeTemplate, KNOWLEDGE_CATEGORIES

try:
    from pptx import Presentation
    from pptx.util import Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.dml import MSO_FILL, MSO_COLOR_TYPE
    from pptx.oxml.ns import qn
    _PPTX_OK = True
    _PPTX_ERR = ''
except Exception as _ex:   # pragma: no cover – Bibliothek fehlt
    _PPTX_OK = False
    _PPTX_ERR = str(_ex)

log = logging.getLogger('pptx_import')

bp = Blueprint('pptx_import', __name__)

EMU_PER_PT = 12700
CANVA_API = 'https://api.canva.com/rest/v1'
CANVA_EXPORT_TIMEOUT_S = 60
CANVA_POLL_INTERVAL_S = 2

PLACEHOLDER_RE = re.compile(r'\{([a-zA-Z_][a-zA-Z0-9_]*)\}')
CITY_ALIASES = {'stadt': 'city_name', 'city': 'city_name', 'stadtname': 'city_name'}

FONT_LABELS = {'anton': 'Anton', 'bold': 'Oswald Bold', 'arial_rounded_bold': 'Arial Rounded Bold'}
FALLBACK_FONT = 'anton'

HELP_TEXT = ('In Canva: Für jede Textstelle, die pro Stadt anders sein soll, schreibst du den '
             'Variablennamen in geschweiften Klammern in die Textbox, zum Beispiel {problem_place} '
             'oder POV: Du musst kurz durch {problem_place}. Fester Text bleibt normal. Dann Teilen → '
             'Herunterladen → PPTX, oder den Design-Link hier einfügen. Mehrere Folien werden zu einer '
             'Serie (Karussell). Gültige Variablen: <Liste>. Eigene Namen sind erlaubt, die KI '
             'formuliert sie dann passend zur Stadt.')

EFFECTS_HINT = 'Konturen/Schatten aus Canva kommen nicht mit, im Studio ergänzen'


# ═══════════════════════════════════════════════════════════════════════════════
# Hilfen: Pfade, Login, Schriften
# ═══════════════════════════════════════════════════════════════════════════════

def _data_root():
    try:
        import app as appmod
        return appmod._DATA_ROOT
    except Exception:
        base = os.path.dirname(os.path.abspath(__file__))
        return os.getenv('MEMEOS_DATA_ROOT') or os.path.join(base, 'instance')


def _dir(sub):
    d = os.path.join(_data_root(), sub)
    os.makedirs(d, exist_ok=True)
    return d


def login_required(f):
    @wraps(f)
    def decorated(*args, **kwargs):
        if not session.get('logged_in'):
            if request.path.startswith('/api/'):
                return jsonify({'error': 'Nicht angemeldet'}), 401
            return redirect('/login')
        return f(*args, **kwargs)
    return decorated


def known_variables():
    """Registry-Schlüssel aus KNOWLEDGE_CATEGORIES plus city_name."""
    keys = [k for k, _label, _color in KNOWLEDGE_CATEGORIES]
    return keys + ['city_name']


def _norm_font(name):
    s = str(name or '').strip().lower()
    s = re.sub(r'\.(ttf|otf|ttc)$', '', s)
    return re.sub(r'[\s\-_]+', '', s)


def font_key(name, bold=False):
    """(key, warnung|None). Nutzt memeos_render.font_key_for, Fallback 'anton'.
    Bei fett wird zuerst die Bold-Variante des Namens probiert."""
    name = (name or '').strip()
    try:
        import memeos_render as mr
    except Exception:
        return FALLBACK_FONT, (f'Schrift {name} nicht vorhanden, Anton verwendet' if name else None)
    if not name:
        return ('bold' if bold else FALLBACK_FONT), None
    candidates = [name]
    if bold and not re.search(r'bold', name, re.I):
        candidates.insert(0, name + ' Bold')
    for cand in candidates:
        try:
            key, exact = mr.font_key_for(cand)
        except Exception:
            continue
        if exact:
            return key, None
    try:
        key, _exact = mr.font_key_for(name)
    except Exception:
        key = FALLBACK_FONT
    known = {_norm_font(k) for k in getattr(mr, 'CANVA_FONT_MAP', {})}
    if _norm_font(name) in known:
        return key, f'Schrift {name} durch {FONT_LABELS.get(key, key)} ersetzt (ähnliche Schrift)'
    return FALLBACK_FONT, f'Schrift {name} nicht vorhanden, Anton verwendet'


# ═══════════════════════════════════════════════════════════════════════════════
# Parser
# ═══════════════════════════════════════════════════════════════════════════════

class _Warnings:
    """Geordnete, duplikatfreie Warnliste."""
    def __init__(self):
        self._items = []
        self._seen = set()

    def add(self, text):
        if text and text not in self._seen:
            self._seen.add(text)
            self._items.append(text)

    def list(self):
        return list(self._items)


def _hex(rgb):
    return '#%02X%02X%02X' % (rgb[0], rgb[1], rgb[2])


def _rgb_tuple(hexstr):
    h = hexstr.lstrip('#')
    return tuple(int(h[i:i + 2], 16) for i in (0, 2, 4))


def _fill_color(fill):
    """('#RRGGBB', alpha 0-1) für eine solide Füllung, sonst (None, reason)."""
    try:
        ftype = fill.type
    except Exception:
        return None, 'unlesbar'
    if ftype is None or ftype == MSO_FILL.BACKGROUND:
        return None, 'keine'
    if ftype != MSO_FILL.SOLID:
        return None, {MSO_FILL.GRADIENT: 'Verlauf', MSO_FILL.PICTURE: 'Bildfüllung',
                      MSO_FILL.PATTERNED: 'Muster', MSO_FILL.TEXTURED: 'Textur'}.get(ftype, 'Sonderfüllung')
    try:
        color = fill.fore_color
        if color.type == MSO_COLOR_TYPE.RGB:
            hexval = _hex(color.rgb)
        else:
            return None, 'Themenfarbe'
    except Exception:
        return None, 'Themenfarbe'
    alpha = 1.0
    try:
        el = fill._xPr if hasattr(fill, '_xPr') else None
        if el is not None:
            a = el.find('.//' + qn('a:solidFill') + '/*/' + qn('a:alpha'))
            if a is not None and a.get('val'):
                alpha = max(0.0, min(1.0, int(a.get('val')) / 100000.0))
    except Exception:
        pass
    return hexval, alpha


def _slide_bg_color(slide):
    """Hintergrundfarbe der Folie ('#RRGGBB') oder None."""
    try:
        if slide.follow_master_background:
            return None
        fill = slide.background.fill
        hexval, _ = _fill_color(fill)
        return hexval
    except Exception:
        return None


def _slide_bg_picture(slide):
    """Bild-Blob eines Folienhintergrunds (bgPr/blipFill) oder None."""
    try:
        bg = slide._element.cSld.bg
        if bg is None:
            return None
        blip = bg.find('.//' + qn('a:blip'))
        if blip is None:
            return None
        rid = blip.get(qn('r:embed'))
        if not rid:
            return None
        return slide.part.related_part(rid).blob
    except Exception:
        return None


def _decode_image(blob):
    im = Image.open(io.BytesIO(blob))
    im.load()
    return im.convert('RGBA')


def _box_brightness(canvas, box):
    """Mittlere Helligkeit (0-255) des Hintergrunds unter einer Box."""
    x, y, w, h = box
    W, H = canvas.size
    x0, y0 = max(0, int(x)), max(0, int(y))
    x1, y1 = min(W, int(x + w)), min(H, int(y + h))
    if x1 <= x0 or y1 <= y0:
        return 255
    region = canvas.crop((x0, y0, x1, y1)).convert('L')
    hist = region.histogram()
    total = sum(hist) or 1
    return sum(i * n for i, n in enumerate(hist)) / total


class _Transform:
    """EMU → Canvas-Pixel, inklusive Gruppen-Koordinatenräume."""
    def __init__(self, scale, ox=0.0, oy=0.0, sx=1.0, sy=1.0):
        self.scale, self.ox, self.oy, self.sx, self.sy = scale, ox, oy, sx, sy

    def box(self, shape):
        left = self.ox + (shape.left or 0) * self.sx
        top = self.oy + (shape.top or 0) * self.sy
        w = (shape.width or 0) * self.sx
        h = (shape.height or 0) * self.sy
        return (left * self.scale, top * self.scale, w * self.scale, h * self.scale)

    def px(self, emu_len):
        return emu_len * self.sx * self.scale

    def child(self, group):
        """Transform für die Kinder einer Gruppe (chOff/chExt → off/ext)."""
        try:
            xfrm = group._element.grpSpPr.xfrm
            off, ext, choff, chext = xfrm.off, xfrm.ext, xfrm.chOff, xfrm.chExt
            if None in (off, ext, choff, chext) or not chext.cx or not chext.cy:
                raise ValueError
            sx = ext.cx / chext.cx
            sy = ext.cy / chext.cy
            ox = self.ox + (off.x - choff.x * sx) * self.sx
            oy = self.oy + (off.y - choff.y * sy) * self.sy
            return _Transform(self.scale, ox, oy, self.sx * sx, self.sy * sy)
        except Exception:
            return _Transform(self.scale, self.ox + (group.left or 0) * self.sx,
                              self.oy + (group.top or 0) * self.sy, self.sx, self.sy)


class _SlideParser:
    def __init__(self, slide, index, canvas_w, canvas_h, scale):
        self.slide = slide
        self.index = index
        self.w, self.h = canvas_w, canvas_h
        self.scale = scale
        self.warn = _Warnings()
        self.elements = []
        self.variables = []
        self.texts = []
        self._text_seen = False
        self._n = 0
        bg_hex = _slide_bg_color(slide)
        self.canvas = Image.new('RGBA', (canvas_w, canvas_h), _rgb_tuple(bg_hex) + (255,) if bg_hex else (255, 255, 255, 255))
        blob = _slide_bg_picture(slide)
        if blob:
            try:
                im = _decode_image(blob).resize((canvas_w, canvas_h), Image.LANCZOS)
                self.canvas.alpha_composite(im)
            except Exception as ex:
                self.warn.add(f'Hintergrundbild der Folie nicht lesbar ({ex.__class__.__name__}), übersprungen')

    def _next_id(self, prefix):
        self._n += 1
        return f's{self.index}_{prefix}{self._n}'

    def _add_var(self, name):
        if name not in self.variables:
            self.variables.append(name)

    # ── Shapes ────────────────────────────────────────────────────────────────
    def walk(self, shapes, tf):
        for shape in shapes:
            try:
                self._shape(shape, tf)
            except Exception as ex:
                log.exception('Shape %s auf Folie %d nicht lesbar', getattr(shape, 'name', '?'), self.index)
                self.warn.add(f'Element "{getattr(shape, "name", "?")}" nicht lesbar ({ex.__class__.__name__}), übersprungen')

    def _shape(self, shape, tf):
        st = shape.shape_type
        if st == MSO_SHAPE_TYPE.GROUP:
            self.walk(shape.shapes, tf.child(shape))
            return
        if st in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE) or hasattr(shape, 'image') and hasattr(shape, 'crop_left'):
            self._picture(shape, tf)
            return
        if st == MSO_SHAPE_TYPE.TABLE or getattr(shape, 'has_table', False):
            self.warn.add(f'Tabelle "{shape.name}" übersprungen')
            return
        if st == MSO_SHAPE_TYPE.CHART or getattr(shape, 'has_chart', False):
            self.warn.add(f'Diagramm "{shape.name}" übersprungen')
            return
        if st in (MSO_SHAPE_TYPE.MEDIA, MSO_SHAPE_TYPE.WEB_VIDEO, MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT,
                  MSO_SHAPE_TYPE.LINKED_OLE_OBJECT, MSO_SHAPE_TYPE.OLE_CONTROL_OBJECT):
            self.warn.add(f'Medien-/Objektelement "{shape.name}" übersprungen')
            return
        if st == MSO_SHAPE_TYPE.LINE:
            self.warn.add(f'Linie "{shape.name}" übersprungen')
            return
        has_text = bool(getattr(shape, 'has_text_frame', False)) and bool(shape.text_frame.text.strip())
        if has_text:
            self._textbox(shape, tf)
        else:
            self._plain_shape(shape, tf)

    def _picture(self, shape, tf):
        try:
            im = _decode_image(shape.image.blob)
        except Exception as ex:
            self.warn.add(f'Bild "{shape.name}" nicht dekodierbar ({ex.__class__.__name__}), übersprungen')
            return
        iw, ih = im.size
        try:
            cl, cr, ct, cb = shape.crop_left or 0, shape.crop_right or 0, shape.crop_top or 0, shape.crop_bottom or 0
        except Exception:
            cl = cr = ct = cb = 0
        if any((cl, cr, ct, cb)):
            box = (int(round(cl * iw)), int(round(ct * ih)), int(round(iw - cr * iw)), int(round(ih - cb * ih)))
            if box[2] > box[0] and box[3] > box[1]:
                im = im.crop(box)
        x, y, w, h = tf.box(shape)
        tw, th = max(1, int(round(w))), max(1, int(round(h)))
        im = im.resize((tw, th), Image.LANCZOS)
        rot = 0.0
        try:
            rot = float(shape.rotation or 0.0)
        except Exception:
            pass
        if abs(rot) > 0.5:
            im = im.rotate(-rot, resample=Image.BICUBIC, expand=True)
            cx, cy = x + w / 2.0, y + h / 2.0
            px, py = int(round(cx - im.width / 2.0)), int(round(cy - im.height / 2.0))
        else:
            px, py = int(round(x)), int(round(y))
        layer = Image.new('RGBA', self.canvas.size, (0, 0, 0, 0))
        layer.paste(im, (px, py), im)
        self.canvas.alpha_composite(layer)
        if self._text_seen:
            self.warn.add(f'Bild "{shape.name}" liegt in Canva über einem Text; hier liegt der Text oben')

    def _shape_geometry(self, shape):
        try:
            return shape.auto_shape_type
        except Exception:
            return None

    def _plain_shape(self, shape, tf):
        """AutoShape ohne Text: solide Füllung wird auf den Hintergrund gezeichnet."""
        try:
            fill = shape.fill
        except Exception:
            return
        hexval, info = _fill_color(fill)
        if not hexval:
            if info == 'keine':
                return
            self.warn.add(f'Form "{shape.name}": Füllung ({info}) nicht übernehmbar, übersprungen')
            return
        alpha = info if isinstance(info, float) else 1.0
        x, y, w, h = tf.box(shape)
        if w < 1 or h < 1:
            return
        try:
            if abs(float(shape.rotation or 0)) > 0.5:
                self.warn.add(f'Form "{shape.name}" ist gedreht; wird ungedreht gezeichnet')
        except Exception:
            pass
        layer = Image.new('RGBA', self.canvas.size, (0, 0, 0, 0))
        d = ImageDraw.Draw(layer)
        color = _rgb_tuple(hexval) + (int(round(alpha * 255)),)
        geom = self._shape_geometry(shape)
        rect = (int(round(x)), int(round(y)), int(round(x + w)) - 1, int(round(y + h)) - 1)
        if geom == MSO_SHAPE.OVAL:
            d.ellipse(rect, fill=color)
        elif geom == MSO_SHAPE.ROUNDED_RECTANGLE:
            d.rounded_rectangle(rect, radius=int(round(self._corner_radius(shape, w, h))), fill=color)
        else:
            d.rectangle(rect, fill=color)
        self.canvas.alpha_composite(layer)

    def _corner_radius(self, shape, w, h):
        adj = 0.16667
        try:
            if len(shape.adjustments):
                adj = float(shape.adjustments[0])
        except Exception:
            pass
        return min(w, h) * adj

    # ── Text ───────────────────────────────────────────────────────────────────
    def _textbox(self, shape, tf):
        tfm = shape.text_frame
        x, y, w, h = tf.box(shape)
        # Füllung der Textbox als rect-Element davor
        try:
            hexval, info = _fill_color(shape.fill)
        except Exception:
            hexval, info = None, 'keine'
        if hexval:
            rect = {'id': self._next_id('r'), 'type': 'rect',
                    'x': round(x), 'y': round(y), 'width': max(1, round(w)), 'height': max(1, round(h)),
                    'fill': hexval}
            if isinstance(info, float) and info < 0.999:
                rect['opacity'] = round(info, 3)
            if self._shape_geometry(shape) == MSO_SHAPE.ROUNDED_RECTANGLE:
                rect['radius'] = round(self._corner_radius(shape, w, h))
            self.elements.append(rect)
        elif info not in ('keine', 'unlesbar'):
            self.warn.add(f'Textbox "{shape.name}": Füllung ({info}) nicht übernehmbar')

        # Innenabstände
        try:
            ml, mr_, mt, mb = (tf.px(tfm.margin_left or 0), tf.px(tfm.margin_right or 0),
                               tf.px(tfm.margin_top or 0), tf.px(tfm.margin_bottom or 0))
            if ml + mr_ < w * 0.8 and mt + mb < h * 0.8:
                x, y, w, h = x + ml, y + mt, w - ml - mr_, h - mt - mb
        except Exception:
            pass

        paragraphs = list(tfm.paragraphs)
        raw_lines = [p.text.replace('\v', '\n').replace('\x0b', '\n') for p in paragraphs]
        text = '\n'.join(raw_lines).strip('\n').strip()
        if not text:
            return
        for alias, target in CITY_ALIASES.items():
            text = re.sub(r'\{' + alias + r'\}', '{' + target + '}', text, flags=re.I)

        runs = [r for p in paragraphs for r in p.runs]
        # Größe: größter Run (pt) → px, Autofit-Skalierung beachten
        size_pt = 0.0
        for r in runs:
            try:
                if r.font.size is not None:
                    size_pt = max(size_pt, float(r.font.size.pt))
            except Exception:
                pass
        if not size_pt:
            for p in paragraphs:
                try:
                    if p.font.size is not None:
                        size_pt = max(size_pt, float(p.font.size.pt))
                except Exception:
                    pass
        if not size_pt:
            size_pt = 18.0
            self.warn.add(f'Textbox "{shape.name}": keine Schriftgröße gefunden, 18 pt angenommen')
        font_scale = 1.0
        try:
            na = tfm._txBody.bodyPr.find(qn('a:normAutofit'))
            if na is not None and na.get('fontScale'):
                font_scale = max(0.1, int(na.get('fontScale')) / 100000.0)
        except Exception:
            pass
        max_size = max(8, int(round(size_pt * EMU_PER_PT * tf.sx * self.scale * font_scale)))
        min_size = max(12, int(round(max_size * 0.5)))
        if min_size > max_size:
            min_size = max_size

        # Schrift + fett + Großschreibung
        bold = False
        font_name = None
        caps_all = False
        biggest = None
        for r in runs:
            try:
                if r.font.bold:
                    bold = True
                if r.font.size is not None and (biggest is None or r.font.size > biggest.font.size):
                    biggest = r
                if r.font._rPr is not None and r.font._rPr.get('cap') == 'all':
                    caps_all = True
            except Exception:
                pass
        for r in ([biggest] if biggest is not None else []) + runs:
            try:
                if r is not None and r.font.name:
                    font_name = r.font.name
                    break
            except Exception:
                pass
        if not font_name:
            try:
                latin = tfm._txBody.find('.//' + qn('a:latin'))
                if latin is not None and latin.get('typeface'):
                    font_name = latin.get('typeface')
            except Exception:
                pass
        key, fwarn = font_key(font_name, bold)
        if fwarn:
            self.warn.add(fwarn)

        # Farbe
        color = None
        for r in ([biggest] if biggest is not None else []) + runs:
            try:
                c = r.font.color
                if c is not None and c.type == MSO_COLOR_TYPE.RGB:
                    color = _hex(c.rgb)
                    break
                if c is not None and c.type is not None:
                    color = 'theme'
            except Exception:
                pass
        if not color or color == 'theme':
            if hexval and (not isinstance(info, float) or info >= 0.5):
                r_, g_, b_ = _rgb_tuple(hexval)
                brightness = 0.299 * r_ + 0.587 * g_ + 0.114 * b_   # Füllung der Textbox zählt
            else:
                brightness = _box_brightness(self.canvas, (x, y, w, h))
            chosen = '#FFFFFF' if brightness < 128 else '#000000'
            if color == 'theme':
                self.warn.add(f'Textbox "{shape.name}": Themenfarbe nicht lesbar, {chosen} nach Hintergrund gewählt')
            color = chosen

        # Ausrichtung
        align = 'left'
        try:
            al = paragraphs[0].alignment
            if al == PP_ALIGN.CENTER:
                align = 'center'
            elif al == PP_ALIGN.RIGHT:
                align = 'right'
        except Exception:
            pass
        valign = 'top'
        try:
            va = tfm.vertical_anchor
            if va == MSO_ANCHOR.MIDDLE:
                valign = 'middle'
            elif va == MSO_ANCHOR.BOTTOM:
                valign = 'bottom'
        except Exception:
            pass
        letters = [ch for ch in text if ch.isalpha()]
        uppercase = caps_all or (bool(letters) and all(ch.isupper() for ch in letters))

        try:
            if abs(float(shape.rotation or 0)) > 0.5:
                self.warn.add(f'Textbox "{shape.name}" ist gedreht; Text wird gerade dargestellt')
        except Exception:
            pass

        el = {'id': self._next_id('t'), 'type': 'text',
              'x': round(x), 'y': round(y), 'width': max(1, round(w)), 'height': max(1, round(h)),
              'font': key, 'max_size': max_size, 'min_size': min_size, 'color': color,
              'align': align, 'valign': valign, 'uppercase': bool(uppercase),
              'line_height': 1.15, 'fit': 'shrink'}
        found = PLACEHOLDER_RE.findall(text)
        whole = PLACEHOLDER_RE.fullmatch(text)
        if whole:
            el['var'] = whole.group(1)
            self._add_var(whole.group(1))
        elif found:
            el['text'] = text
            el['vars'] = list(dict.fromkeys(found))
            for v in found:
                self._add_var(v)
        else:
            el['text'] = text
        self.elements.append(el)
        self.texts.append(text)
        self._text_seen = True
        self.warn.add(EFFECTS_HINT)

    # ── Ergebnis ──────────────────────────────────────────────────────────────
    def result(self):
        return {
            'slide': self.index,
            'canvas': {'width': self.w, 'height': self.h},
            'background': self.canvas,
            'elements': self.elements,
            'variables': list(self.variables),
            'texts': list(self.texts),
            'warnings': self.warn.list(),
        }


def parse_pptx(path, target_width=1080):
    """Liest eine PPTX-Datei; je Folie {'slide','canvas','background' (PIL RGBA),'elements',
    'variables','texts','warnings'}. path darf Pfad oder Dateiobjekt sein."""
    if not _PPTX_OK:
        raise RuntimeError(f'python-pptx nicht installiert ({_PPTX_ERR}) – pip install python-pptx==1.0.2')
    prs = Presentation(path)
    sw, sh = int(prs.slide_width or 0), int(prs.slide_height or 0)
    if sw <= 0 or sh <= 0:
        raise ValueError('Foliengröße nicht lesbar')
    scale = float(target_width) / sw
    canvas_w = int(target_width)
    canvas_h = max(1, int(round(sh * scale)))
    out = []
    for i, slide in enumerate(prs.slides, start=1):
        parser = _SlideParser(slide, i, canvas_w, canvas_h, scale)
        parser.walk(slide.shapes, _Transform(scale))
        out.append(parser.result())
    return out


# ═══════════════════════════════════════════════════════════════════════════════
# Speichern
# ═══════════════════════════════════════════════════════════════════════════════

def _validate(config):
    try:
        import memeos_render as mr
        return list(mr.validate_config(config) or [])
    except ImportError:
        return []
    except Exception as ex:
        return [f'Config-Prüfung fehlgeschlagen: {ex}']


def _element_summary(el):
    s = {'id': el.get('id'), 'type': el.get('type'),
         'box': [el.get('x'), el.get('y'), el.get('width'), el.get('height')]}
    if el.get('type') == 'text':
        s.update({'var': el.get('var'), 'text': el.get('text'), 'font': el.get('font'),
                  'max_size': el.get('max_size'), 'color': el.get('color')})
    elif el.get('type') == 'rect':
        s['fill'] = el.get('fill')
    return s


def _safe_name(s, fallback='Import'):
    s = re.sub(r'[\r\n\t]+', ' ', str(s or '')).strip()
    return s[:180] or fallback


def import_slides(slides, name, category='allgemein', series=None, dry_run=False,
                  source='PPTX', canva_url=''):
    """Folien-Analysen in MemeTemplate-Zeilen überführen (eine Transaktion). Liefert das
    Antwort-Dict der Import-Routen."""
    name = _safe_name(name)
    category = (category or 'allgemein').strip() or 'allgemein'
    multi = len(slides) > 1
    series_name = (series or '').strip() or (name if multi else None)
    global_warnings = _Warnings()
    if not slides:
        global_warnings.add('Die Datei enthält keine Folien')
    results = []
    configs = []
    for s in slides:
        config = {'canvas': s['canvas'], 'elements': s['elements']}
        errors = _validate(config)
        warns = list(s['warnings'])
        for e in errors:
            warns.append(f'Config-Prüfung: {e}')
        if not s['elements']:
            warns.append('Keine Textboxen gefunden – Template hat nur ein Hintergrundbild')
        tname = f'{name} – Folie {s["slide"]}' if multi else name
        results.append({
            'id': None, 'name': tname, 'slide': s['slide'], 'canvas': s['canvas'],
            'elements_count': len(s['elements']), 'variables': s['variables'],
            'unknown_variables': [v for v in s['variables'] if v not in known_variables()],
            'elements': [_element_summary(e) for e in s['elements']],
            'warnings': warns, 'preview_url': '', 'studio_url': '',
        })
        configs.append(config)

    if dry_run:
        return {'templates': results, 'warnings': global_warnings.list(), 'dry_run': True,
                'series': series_name}

    upload_dir = _dir('uploads')
    ts = datetime.now().strftime('%Y%m%d_%H%M%S_%f')[:-3]
    written = []
    try:
        import app as appmod
    except Exception:
        appmod = None
    try:
        for s, res, config in zip(slides, results, configs):
            fname = f'import_{ts}_{s["slide"]}.png'
            local = os.path.join(upload_dir, fname)
            s['background'].convert('RGB').save(local, format='PNG', optimize=True)
            written.append(local)
            preview_url = ''
            if appmod is not None and hasattr(appmod, '_upload_cloudinary'):
                try:
                    preview_url = appmod._upload_cloudinary(local, folder='memeos/templates',
                                                            resource_type='image') or ''
                except Exception as ex:
                    log.warning('Cloudinary-Upload für %s fehlgeschlagen: %s', fname, ex)
            notes = '\n'.join(res['warnings'])
            t = MemeTemplate(
                name=res['name'],
                description=f'Importiert aus {source} am {datetime.now().strftime("%d.%m.%Y")}',
                canva_url=canva_url or '',
                render_type='pil',
                pil_config=json.dumps(config, ensure_ascii=False),
                required_vars=json.dumps(s['variables']),
                tags=json.dumps(['import']),
                category=category,
                preview_image=fname,
                preview_url=preview_url or None,
                example_text=(' | '.join(s['texts']))[:300],
                notes=notes,
                series=series_name,
                series_position=(s['slide'] if series_name else None),
            )
            db.session.add(t)
            db.session.flush()
            res['id'] = t.id
            res['preview_url'] = preview_url or f'/uploads/{fname}'
            res['studio_url'] = f'/studio/{t.id}'
        db.session.commit()
    except Exception as ex:
        db.session.rollback()
        for p in written:
            try:
                os.remove(p)
            except OSError:
                pass
        log.exception('Template-Import fehlgeschlagen')
        raise RuntimeError(f'Speichern fehlgeschlagen: {ex}') from ex
    return {'templates': results, 'warnings': global_warnings.list(), 'dry_run': False,
            'series': series_name}


def _truthy(v):
    return str(v or '').strip().lower() in ('1', 'true', 'yes', 'ja', 'on')


# ═══════════════════════════════════════════════════════════════════════════════
# Canva
# ═══════════════════════════════════════════════════════════════════════════════

_DESIGN_ID_RE = re.compile(r'canva\.com/design/([A-Za-z0-9_-]+)')


def canva_design_id(url, resolve_short_links=True):
    """Design-ID aus Canva-URL (oder canva.link-Kurzlink, wird per Redirect aufgelöst)."""
    s = (url or '').strip()
    if not s:
        return None
    m = _DESIGN_ID_RE.search(s)
    if m:
        return m.group(1)
    if re.fullmatch(r'[A-Za-z0-9_-]{8,}', s):
        return s
    if resolve_short_links and re.search(r'https?://(www\.)?canva\.link/', s):
        try:
            r = requests.get(s, allow_redirects=True, timeout=15,
                             headers={'User-Agent': 'MemeOS/1.0'})
            m = _DESIGN_ID_RE.search(r.url or '')
            if m:
                return m.group(1)
        except Exception as ex:
            log.warning('canva.link nicht auflösbar: %s', ex)
    return None


def _canva_error(resp, fallback):
    if resp.status_code == 403:
        return 'Kein Zugriff auf dieses Design'
    if resp.status_code == 404:
        return 'Design nicht gefunden – Link prüfen'
    if resp.status_code == 401:
        return 'Canva-Anmeldung abgelaufen. Einstellungen → Canva neu verbinden'
    if resp.status_code == 429:
        return 'Canva-Limit erreicht, bitte kurz warten'
    try:
        msg = resp.json().get('message') or resp.json().get('error', {}).get('message')
    except Exception:
        msg = None
    return f'{fallback} ({resp.status_code}{": " + msg if msg else ""})'


def canva_export_pptx(design_id, token, dest_dir, timeout_s=CANVA_EXPORT_TIMEOUT_S):
    """Startet den PPTX-Export, pollt bis Erfolg, lädt herunter. Liefert (pfad, None) oder
    (None, fehlertext)."""
    headers = {'Authorization': f'Bearer {token}', 'Content-Type': 'application/json'}
    try:
        r = requests.post(f'{CANVA_API}/exports',
                          json={'design_id': design_id, 'format': {'type': 'pptx'}},
                          headers=headers, timeout=20)
    except requests.RequestException as ex:
        return None, f'Canva nicht erreichbar: {ex.__class__.__name__}'
    if not r.ok:
        return None, _canva_error(r, 'Export konnte nicht gestartet werden')
    try:
        job = r.json().get('job') or {}
    except Exception:
        job = {}
    job_id = job.get('id')
    if not job_id:
        return None, 'Canva hat keine Export-Job-ID geliefert'
    deadline = time.time() + timeout_s
    status = job.get('status')
    while status not in ('success', 'failed'):
        if time.time() >= deadline:
            return None, f'Canva-Export nicht in {timeout_s} s fertig – später erneut versuchen'
        time.sleep(CANVA_POLL_INTERVAL_S)
        try:
            pr = requests.get(f'{CANVA_API}/exports/{job_id}', headers=headers, timeout=20)
        except requests.RequestException as ex:
            return None, f'Canva nicht erreichbar: {ex.__class__.__name__}'
        if not pr.ok:
            return None, _canva_error(pr, 'Export-Status nicht abrufbar')
        try:
            job = pr.json().get('job') or {}
        except Exception:
            job = {}
        status = job.get('status')
    if status == 'failed':
        err = (job.get('error') or {})
        return None, f'Canva-Export fehlgeschlagen: {err.get("message") or err.get("code") or "unbekannt"}'
    urls = job.get('urls') or (job.get('result') or {}).get('urls') or []
    if not urls:
        return None, 'Canva-Export ohne Download-Link'
    os.makedirs(dest_dir, exist_ok=True)
    path = os.path.join(dest_dir, f'canva_{re.sub(r"[^A-Za-z0-9_-]", "", design_id)}_{datetime.now().strftime("%Y%m%d_%H%M%S")}.pptx')
    try:
        with requests.get(urls[0], stream=True, timeout=60) as dl:
            if not dl.ok:
                return None, f'Download der PPTX fehlgeschlagen ({dl.status_code})'
            with open(path, 'wb') as fh:
                for chunk in dl.iter_content(1 << 16):
                    if chunk:
                        fh.write(chunk)
    except requests.RequestException as ex:
        return None, f'Download der PPTX fehlgeschlagen: {ex.__class__.__name__}'
    return path, None


# ═══════════════════════════════════════════════════════════════════════════════
# Routen
# ═══════════════════════════════════════════════════════════════════════════════

@bp.route('/api/templates/import/help', methods=['GET'])
@login_required
def api_import_help():
    variables = known_variables()
    labels = {k: re.sub(r'^[^\w]+', '', label).strip() for k, label, _c in KNOWLEDGE_CATEGORIES}
    labels['city_name'] = 'Stadtname'
    return jsonify({
        'anleitung': HELP_TEXT.replace('<Liste>', ', '.join('{' + v + '}' for v in variables)),
        'variables': variables,
        'labels': labels,
        'aliases': {'{' + a + '}': '{' + t + '}' for a, t in CITY_ALIASES.items()},
        'formate': ['pptx'],
        'canva_verbunden': _canva_connected(),
        'hinweise': [
            EFFECTS_HINT,
            'Bilder und Farbflächen werden zu einem Hintergrundbild zusammengefügt, Textboxen bleiben bearbeitbar',
            'Tabellen, Diagramme, Videos und Linien werden übersprungen',
            'Gedrehte Elemente werden gerade dargestellt',
        ],
    })


def _canva_connected():
    try:
        import app as appmod
        return bool(appmod._canva_is_connected())
    except Exception:
        return False


def _import_error(ex):
    msg = str(ex)
    if isinstance(ex, RuntimeError) and 'python-pptx' in msg:
        return jsonify({'error': msg}), 500
    return jsonify({'error': f'PPTX nicht lesbar: {msg}'}), 400


@bp.route('/api/templates/import/pptx', methods=['POST'])
@login_required
def api_import_pptx():
    f = request.files.get('file')
    if f is None or not f.filename:
        return jsonify({'error': 'Keine Datei übergeben (Feld "file")'}), 400
    if not f.filename.lower().endswith('.pptx'):
        return jsonify({'error': 'Nur .pptx-Dateien werden unterstützt'}), 400
    form = request.form
    dry_run = _truthy(form.get('dry_run'))
    name = (form.get('name') or '').strip() or os.path.splitext(os.path.basename(f.filename))[0]
    category = (form.get('category') or 'allgemein').strip() or 'allgemein'
    series = (form.get('series') or '').strip() or None

    data = f.read()
    if not data:
        return jsonify({'error': 'Datei ist leer'}), 400
    stored = None
    if not dry_run:
        imports_dir = _dir('imports')
        safe = re.sub(r'[^A-Za-z0-9._-]+', '_', os.path.basename(f.filename))[:80]
        stored = os.path.join(imports_dir, f'upload_{datetime.now().strftime("%Y%m%d_%H%M%S")}_{safe}')
        with open(stored, 'wb') as fh:
            fh.write(data)
    try:
        slides = parse_pptx(io.BytesIO(data))
    except Exception as ex:
        log.warning('PPTX-Parse fehlgeschlagen: %s', ex)
        if stored:
            try:
                os.remove(stored)
            except OSError:
                pass
        return _import_error(ex)
    try:
        result = import_slides(slides, name, category, series, dry_run, source='PPTX')
    except RuntimeError as ex:
        return jsonify({'error': str(ex)}), 500
    result['source'] = 'pptx'
    result['file'] = os.path.basename(stored) if stored else None
    return jsonify(result)


@bp.route('/api/templates/import/canva', methods=['POST'])
@login_required
def api_import_canva():
    d = request.get_json(silent=True) or {}
    design_url = (d.get('design_url') or d.get('url') or '').strip()
    if not design_url:
        return jsonify({'error': 'design_url fehlt'}), 400
    try:
        import app as appmod
        token = appmod._canva_get_token()
    except Exception as ex:
        log.warning('Canva-Token nicht abrufbar: %s', ex)
        token = None
    if not token:
        return jsonify({'error': 'Canva nicht verbunden. Einstellungen → Canva verbinden'}), 400
    design_id = canva_design_id(design_url)
    if not design_id:
        return jsonify({'error': 'Keine Design-ID im Link gefunden. Erwartet: https://www.canva.com/design/<ID>/...'}), 400
    dry_run = _truthy(d.get('dry_run'))
    name = (d.get('name') or '').strip() or f'Canva {design_id}'
    category = (d.get('category') or 'allgemein').strip() or 'allgemein'
    series = (d.get('series') or '').strip() or None

    path, err = canva_export_pptx(design_id, token, _dir('imports'))
    if err:
        status = 403 if err.startswith('Kein Zugriff') else (404 if err.startswith('Design nicht gefunden') else 502)
        return jsonify({'error': err, 'design_id': design_id}), status
    try:
        slides = parse_pptx(path)
    except Exception as ex:
        log.warning('Canva-PPTX-Parse fehlgeschlagen: %s', ex)
        return _import_error(ex)
    try:
        result = import_slides(slides, name, category, series, dry_run, source='Canva',
                               canva_url=f'https://www.canva.com/design/{design_id}/')
    except RuntimeError as ex:
        return jsonify({'error': str(ex)}), 500
    if dry_run:
        try:
            os.remove(path)
        except OSError:
            pass
    result['source'] = 'canva'
    result['design_id'] = design_id
    result['file'] = None if dry_run else os.path.basename(path)
    return jsonify(result)


def init_app(flask_app):
    if 'pptx_import' in flask_app.blueprints:
        return
    flask_app.register_blueprint(bp)
    if not _PPTX_OK:
        log.warning('python-pptx fehlt (%s) – Template-Import liefert 500 bis zur Installation', _PPTX_ERR)
